"""
Text extraction utilities for different file types.

Supports:
- PDF: Using PyPDF2
- DOCX: Using python-docx
- PPTX: Using python-pptx
- TXT: Direct read

Each extractor handles its file type independently.
"""

import io
from abc import ABC, abstractmethod
from pathlib import Path
from typing import BinaryIO

from app.core.logging import get_logger

logger = get_logger(__name__)


class TextExtractor(ABC):
    """Abstract base class for text extractors."""
    
    @abstractmethod
    def extract(self, file_content: BinaryIO) -> str:
        """
        Extract text from file content.
        
        Args:
            file_content: Binary file-like object.
            
        Returns:
            Extracted text as string.
        """
        pass


class PDFExtractor(TextExtractor):
    """Extract text from PDF files using PyPDF2."""
    
    def extract(self, file_content: BinaryIO) -> str:
        """Extract text from all pages of a PDF."""
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(file_content)
            text_parts = []
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num}: {e}")
                    continue
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise ExtractionError(f"Failed to extract text from PDF: {e}")


class DOCXExtractor(TextExtractor):
    """Extract text from DOCX files using python-docx."""
    
    def extract(self, file_content: BinaryIO) -> str:
        """Extract text from all paragraphs in a DOCX."""
        try:
            from docx import Document
            
            doc = Document(file_content)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise ExtractionError(f"Failed to extract text from DOCX: {e}")


class PPTXExtractor(TextExtractor):
    """Extract text from PPTX files using python-pptx."""
    
    def extract(self, file_content: BinaryIO) -> str:
        """Extract text from all slides in a PPTX."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_content)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise ExtractionError(f"Failed to extract text from PPTX: {e}")


class TXTExtractor(TextExtractor):
    """Extract text from plain text files."""
    
    def extract(self, file_content: BinaryIO) -> str:
        """Read text directly from file."""
        try:
            # Try UTF-8 first, fall back to latin-1
            content = file_content.read()
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1")
        except Exception as e:
            logger.error(f"TXT extraction failed: {e}")
            raise ExtractionError(f"Failed to read text file: {e}")


class ExtractionError(Exception):
    """Custom exception for text extraction failures."""
    pass


# Mapping of file types to extractors
EXTRACTORS: dict[str, type[TextExtractor]] = {
    "pdf": PDFExtractor,
    "docx": DOCXExtractor,
    "pptx": PPTXExtractor,
    "ppt": PPTXExtractor,  # Try to handle .ppt as .pptx (may not work for old format)
    "txt": TXTExtractor,
    "md": TXTExtractor,
}

# Supported file extensions
SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "ppt", "txt", "md"}


def get_file_extension(filename: str) -> str:
    """
    Get lowercase file extension without dot.
    
    Args:
        filename: Original filename.
        
    Returns:
        Lowercase extension (e.g., 'pdf', 'docx').
    """
    return Path(filename).suffix.lower().lstrip(".")


def is_supported_file(filename: str) -> bool:
    """
    Check if file type is supported.
    
    Args:
        filename: Original filename.
        
    Returns:
        True if file type is supported.
    """
    ext = get_file_extension(filename)
    return ext in SUPPORTED_EXTENSIONS


def extract_text(file_content: bytes, filename: str) -> str:
    """
    Extract text from file based on its type.
    
    Args:
        file_content: Raw file bytes.
        filename: Original filename (used to determine type).
        
    Returns:
        Extracted text as string.
        
    Raises:
        ValueError: If file type is not supported.
        ExtractionError: If extraction fails.
    """
    ext = get_file_extension(filename)
    
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file type: {ext}")
    
    extractor_class = EXTRACTORS[ext]
    extractor = extractor_class()
    
    # Wrap bytes in BytesIO for file-like interface
    file_io = io.BytesIO(file_content)
    
    return extractor.extract(file_io)
